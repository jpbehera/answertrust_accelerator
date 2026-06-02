"""Convert 04_MEGA_AnswerTrust_deck.html to a PPTX (one PNG per slide stage).

Slides with click-driven stages are expanded so every reveal state becomes its
own PPTX slide:
  - #slide-1b-anim  : data-step 1 → 2 → 3
  - #slide-stack    : data-step 1 → 2
  - #slide-value    : data-stage 0 → 1 → 2
"""
import asyncio
from pathlib import Path
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

HERE = Path(__file__).parent
HTML = HERE / "04_MEGA_AnswerTrust_deck.html"
OUT_DIR = HERE / "_pptx_build"
PPTX = HERE / "04_MEGA_AnswerTrust_deck.pptx"

SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
DPI = 150  # render scale
W_PX, H_PX = int(SLIDE_W_IN * DPI), int(SLIDE_H_IN * DPI)

# slide_id -> (attribute_name, [values to render in order])
STAGE_MAP = {
    "slide-1b-anim": ("data-step",  ["1", "2", "3"]),
    "slide-stack":   ("data-step",  ["1", "2"]),
    "slide-value":   ("data-stage", ["0", "1", "2"]),
}

async def render():
    OUT_DIR.mkdir(exist_ok=True)
    # wipe stale frames so removed slides don't leak into the deck
    for old in OUT_DIR.glob("slide_*.png"):
        old.unlink()

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(
            viewport={"width": W_PX, "height": H_PX},
            device_scale_factor=2,
        )
        page = await ctx.new_page()
        await page.goto(HTML.as_uri(), wait_until="networkidle")
        # Strip browser-only chrome (page margins) for clean capture, and
        # disable entrance animations so delayed reveals (e.g. welcome-slide
        # pillars) are captured in their final state.
        await page.add_style_tag(content="""
            body { margin:0 !important; padding:0 !important; background:#fff !important; }
            .slide { margin:0 !important; box-shadow:none !important; page-break-after:auto !important; }
            *, *::before, *::after {
              animation-duration: 0.001s !important;
              animation-delay: 0s !important;
              animation-iteration-count: 1 !important;
              transition-duration: 0.001s !important;
              transition-delay: 0s !important;
            }
        """)

        slides = await page.query_selector_all("section.slide")
        paths = []
        frame_idx = 0
        for el in slides:
            sid = await el.get_attribute("id") or ""
            stages = STAGE_MAP.get(sid)

            if stages:
                attr, values = stages
                for v in values:
                    # set the stage attribute and let any CSS transitions settle
                    await page.evaluate(
                        "([id, a, val]) => document.getElementById(id).setAttribute(a, val)",
                        [sid, attr, v],
                    )
                    # On the final stage of slide-1b-anim, also reveal the
                    # SVG pillars + flow lines (normally gated by clicking
                    # the central AnswerTrust hub).
                    if sid == "slide-1b-anim" and v == values[-1]:
                        await page.evaluate(
                            "document.querySelectorAll('#slide-1b-anim .tissue-svg')"
                            ".forEach(s => s.classList.add('revealed'))"
                        )
                    await page.wait_for_timeout(900)
                    frame_idx += 1
                    fp = OUT_DIR / f"slide_{frame_idx:02d}_{sid}_{attr}-{v}.png"
                    await el.screenshot(path=str(fp))
                    paths.append(fp)
                    print(f"  rendered {fp.name}")
            else:
                frame_idx += 1
                fp = OUT_DIR / f"slide_{frame_idx:02d}_{sid or 'slide'}.png"
                await el.screenshot(path=str(fp))
                paths.append(fp)
                print(f"  rendered {fp.name}")

        await browser.close()
        return paths

def build_pptx(images):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    for img in images:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(PPTX)
    print(f"  wrote {PPTX.name} ({len(images)} slides)")

async def main():
    print("rendering slides...")
    imgs = await render()
    print("building pptx...")
    build_pptx(imgs)

if __name__ == "__main__":
    asyncio.run(main())
